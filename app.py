from flask import Flask, request, jsonify, send_file
from pptx import Presentation
import os
import logging
import sys

app = Flask(__name__)

# ログ設定
logging.basicConfig(stream=sys.stdout, level=logging.DEBUG)
app.logger.setLevel(logging.DEBUG)

# 許可されたAPIキー
API_KEY = "MySecureAPIKey123"

# APIキー認証デコレータ
def require_api_key(func):
    def wrapper(*args, **kwargs):
        api_key = request.headers.get("x_api_key")
        if not api_key:
            app.logger.warning("Missing API key in headers.")
            return jsonify({"error": "Unauthorized: API key missing"}), 401
        if api_key != API_KEY:
            app.logger.warning(f"Invalid API key: {api_key}")
            return jsonify({"error": "Unauthorized: Invalid API key"}), 401
        return func(*args, **kwargs)
    wrapper.__name__ = func.__name__
    return wrapper

# リクエストヘッダーをログに記録
@app.before_request
def log_request_headers():
    app.logger.info(f"Request headers: {dict(request.headers)}")

# PPTファイルの保存パスを生成
def get_presentation_path(presentation_id):
    return f"./presentations/{presentation_id}.pptx"

# /create-slide エンドポイント
@app.route('/create-slide', methods=['POST'])
@require_api_key
def create_slide():
    app.logger.debug("Received request at /create-slide endpoint")
    try:
        data = request.json
        app.logger.debug(f"Request payload: {data}")

        title = data.get('title')
        content = data.get('content')
        presentation_id = data.get('presentationId')
        slide_layout = data.get('slideLayout', 'Title and Content')

        if not title or not content or not presentation_id:
            app.logger.error("Missing required fields in request")
            return jsonify({"error": "Missing required fields"}), 400

        file_path = get_presentation_path(presentation_id)
        app.logger.debug(f"Target file path for presentation: {file_path}")

        if not os.path.exists(file_path):
            os.makedirs(os.path.dirname(file_path), exist_ok=True)
            app.logger.info(f"Creating new PowerPoint file at {file_path}")
            presentation = Presentation()
            presentation.save(file_path)

        app.logger.debug(f"Opening PowerPoint file at {file_path}")
        presentation = Presentation(file_path)
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        presentation.save(file_path)
        app.logger.info(f"Slide '{title}' successfully added to {presentation_id}")

        return jsonify({
            "success": True,
            "message": f"Slide '{title}' added to presentation '{presentation_id}' successfully.",
            "presentationId": presentation_id
        })
    except Exception as e:
        app.logger.exception("An error occurred while creating a slide")
        return jsonify({"error": f"Failed to create slide: {str(e)}"}), 500

# /download-presentation エンドポイント
@app.route('/download-presentation', methods=['GET'])
@require_api_key
def download_presentation():
    app.logger.debug("Received request at /download-presentation endpoint")
    try:
        presentation_id = request.args.get('presentationId')
        app.logger.debug(f"Received presentationId: {presentation_id}")

        if not presentation_id:
            app.logger.error("No presentationId provided")
            return jsonify({"error": "presentationId is required"}), 400

        file_path = get_presentation_path(presentation_id)
        app.logger.debug(f"Looking for file at: {file_path}")

        if not os.path.exists(file_path):
            app.logger.error(f"File not found: {file_path}")
            return jsonify({"error": "Presentation not found"}), 404

        app.logger.info(f"Serving file: {file_path}")
        return send_file(file_path, as_attachment=True)
    except Exception as e:
        app.logger.exception("An error occurred while serving the file")
        return jsonify({"error": f"Failed to download presentation: {str(e)}"}), 500

# ルートエンドポイント
@app.route('/', methods=['GET'])
def root():
    return jsonify({"message": "Welcome to the PowerPoint API!"})

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 10000))
    app.run(debug=True, host='0.0.0.0', port=port)
